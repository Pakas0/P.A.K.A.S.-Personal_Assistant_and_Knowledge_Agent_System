from pptx import Presentation
from pptx.util import Inches

def build_pptx(title: str, slides: list[dict], output_path: str) -> str:
    """
    slides: list of {"title": str, "content": str} — 1 slide per item,
    layout sederhana (title + content placeholder), TIDAK ada styling kompleks/gambar.
    """
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title

    content_layout = prs.slide_layouts[1]
    for s in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = s.get("heading") or s.get("title", "")
        body = slide.placeholders[1].text_frame
        body.text = s.get("content", "")
    prs.save(output_path)
    return output_path
